#!/usr/bin/env python3
"""competitors.md 내용 -> 경쟁사 분석 PPT.
표지 / 시장개요 / 비교표 / 차별화 / 추천 슬라이드.
실행: python3 make_competitor_pptx.py  ->  out/competitor_analysis.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.join(HERE, "out")
os.makedirs(OUT_DIR, exist_ok=True)
OUT = os.path.join(OUT_DIR, "competitor_analysis.pptx")

NAVY = RGBColor(0x2F, 0x54, 0x96)
GRAY = RGBColor(0x59, 0x59, 0x59)
RED = RGBColor(0xC0, 0x39, 0x39)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_text(slide, l, t, w, h, text, size=18, bold=False, color=GRAY,
             align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    return tb


def add_bullets(slide, l, t, w, h, items, size=16):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - lvl * 2)
        r.font.color.rgb = NAVY if lvl == 0 else GRAY
        r.font.bold = (lvl == 0)
        p.space_after = Pt(6)
    return tb


def band(slide, color=NAVY, h=Inches(1.1)):
    bar = slide.shapes.add_shape(1, 0, 0, SW, h)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def header(slide, title, idx):
    band(slide)
    add_text(slide, Inches(0.5), Inches(0.22), Inches(11), Inches(0.7),
             title, size=26, bold=True, color=WHITE)
    add_text(slide, Inches(12.3), Inches(0.35), Inches(0.8), Inches(0.5),
             str(idx), size=18, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


# ===== 1. 표지 =====
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
add_text(s, Inches(1), Inches(2.4), Inches(11.3), Inches(1.2),
         "성수동 디저트 카페 경쟁 분석", size=44, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.8),
         "데일리브루(자사) 포지셔닝 & 차별화 전략", size=22, color=RGBColor(0xCF, 0xDA, 0xF0),
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.5),
         "AI 공장장 부트캠프 · 리뷰/경쟁 리포트 자동화", size=14,
         color=RGBColor(0xAD, 0xBE, 0xDB), align=PP_ALIGN.CENTER)

# ===== 2. 시장 개요 =====
s = prs.slides.add_slide(BLANK)
header(s, "1. 시장 개요", 2)
add_bullets(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5),
            [("성수동 = 카페 밀집 상권, 20~30대 여성·SNS 인증 수요가 핵심", 0),
             ("디저트(치즈케이크·크로플) 중심 객단가 상승, 주말 웨이팅이 일상", 0),
             ("자사 VoC 분석: 부정 리뷰 1위 = '대기시간' (주말 피크 집중)", 0),
             ("→ 회전율·예약 동선이 핵심 경쟁 변수", 1),
             ("디저트 품질·친절은 호평 — 운영(대기·좌석)에서 이탈 발생", 0)], size=20)

# ===== 3. 경쟁사 비교표 =====
s = prs.slides.add_slide(BLANK)
header(s, "2. 경쟁사 비교표", 3)
table_data = [
    ["항목", "데일리브루(자사)", "어반로스터스", "슈가힐", "모닝릿"],
    ["콘셉트", "디저트 강점 동네카페", "스페셜티 로스팅", "비주얼 디저트", "브런치+카페"],
    ["시그니처", "치즈케이크·크로플", "싱글오리진 드립", "시즌 케이크", "브런치 플레이트"],
    ["객단가", "9,000원", "12,000원", "11,000원", "18,000원"],
    ["좌석", "24석(협소)", "40석", "30석", "55석"],
    ["주말 웨이팅", "길다(이슈)", "보통", "매우 길다", "보통"],
    ["약점", "좌석·대기", "디저트 약함", "가격·대기", "커피 평범"],
]
rows, cols = len(table_data), len(table_data[0])
gt = s.shapes.add_table(rows, cols, Inches(0.5), Inches(1.4),
                        Inches(12.3), Inches(5.4)).table
gt.columns[0].width = Inches(2.0)
for c in range(1, cols):
    gt.columns[c].width = Inches(2.575)
for r in range(rows):
    for c in range(cols):
        cell = gt.cell(r, c)
        cell.text = table_data[r][c]
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(13)
        if r == 0:
            para.font.bold = True; para.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        elif c == 1:
            para.font.bold = True; para.font.color.rgb = NAVY
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xDE, 0xE7, 0xF5)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT

# ===== 4. 차별화 전략 =====
s = prs.slides.add_slide(BLANK)
header(s, "3. 자사 차별화 전략", 4)
add_bullets(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5),
            [("'맛은 슈가힐급, 가격은 합리적' 포지션", 0),
             ("디저트 품질 호평을 핵심 메시지로 전면화", 1),
             ("대기시간 약점 정면 대응", 0),
             ("주말 모바일 예약/픽업 주문, 좌석 회전 동선 개선 (VoC 1위 직접 해소)", 1),
             ("소규모(24석)의 강점 = 큐레이션", 0),
             ("시즌 한정 디저트로 '여기서만' 경험 제공 → 재방문 동기", 1)], size=20)

# ===== 5. 추천 액션 =====
s = prs.slides.add_slide(BLANK)
header(s, "4. 추천 액션", 5)
box = s.shapes.add_shape(1, Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.2))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFC, 0xE4, 0xE4)
box.line.color.rgb = RED
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "⭐ 핵심 액션: 로컬 마이크로 인플루언서 협업 (→ 프로젝트 9 연계)"
r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = RED
add_bullets(s, Inches(0.7), Inches(3.0), Inches(11.9), Inches(4),
            [("SNS 노출 약점을 마이크로 인플루언서로 보완 — '맛+가성비' 인증샷 확산", 0),
             ("주말 예약/픽업 도입 → 대기 불만 완화 → 재방문율·평점 회복", 0),
             ("시즌 한정 디저트 캘린더 운영 → 재방문 동기 부여", 0)], size=19)

prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
